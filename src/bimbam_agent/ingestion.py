"""Leitores modulares para os documentos usados pelo pipeline RAG.

Cada leitor converte um arquivo em :class:`DocumentFragment` sem misturar
localizações lógicas. Assim, páginas, slides, planilhas e seções continuam
rastreáveis durante a recuperação e a geração de citações.
"""

from __future__ import annotations

import csv
import json
import re
from collections.abc import Callable, Iterable, Mapping
from html import unescape
from io import StringIO
from pathlib import Path
from typing import Any

from .models import DocumentFragment


class IngestionError(RuntimeError):
    """Falha ao localizar, interpretar ou extrair um documento."""


class UnsupportedDocumentError(IngestionError):
    """O formato do documento não possui leitor registrado."""


class EmptyDocumentError(IngestionError):
    """O documento não possui conteúdo textual utilizável."""


Reader = Callable[[Path, Mapping[str, Any]], list[DocumentFragment]]

SUPPORTED_EXTENSIONS = frozenset(
    {
        ".pdf",
        ".docx",
        ".xlsx",
        ".pptx",
        ".md",
        ".markdown",
        ".txt",
        ".csv",
        ".json",
        ".html",
        ".htm",
    }
)


def clean_text(value: str) -> str:
    """Normaliza ruído comum sem remover listas, títulos ou quebras semânticas."""

    text = unescape(value).replace("\ufeff", "").replace("\u00a0", " ")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[\t ]+", " ", line).strip() for line in text.split("\n")]

    cleaned: list[str] = []
    blank = False
    for line in lines:
        if line:
            cleaned.append(line)
            blank = False
        elif cleaned and not blank:
            cleaned.append("")
            blank = True
    return "\n".join(cleaned).strip()


def load_manifest(path: str | Path | None) -> dict[str, dict[str, Any]]:
    """Carrega um manifesto opcional no formato ``{"documents": {...}}``.

    Também é aceito um objeto simples que mapeie diretamente nomes de arquivo
    para metadados, o que facilita o uso da função fora do projeto.
    """

    if path is None:
        return {}

    manifest_path = Path(path)
    if not manifest_path.is_file():
        raise IngestionError(f"Manifesto não encontrado: {manifest_path}")

    try:
        payload = json.loads(_read_text(manifest_path))
    except (OSError, UnicodeError, json.JSONDecodeError) as exc:
        raise IngestionError(f"Não foi possível ler o manifesto '{manifest_path}': {exc}") from exc

    if not isinstance(payload, dict):
        raise IngestionError(
            f"Manifesto inválido em '{manifest_path}': a raiz deve ser um objeto JSON"
        )

    documents = payload.get("documents", payload)
    if not isinstance(documents, dict):
        raise IngestionError(
            f"Manifesto inválido em '{manifest_path}': 'documents' deve ser um objeto"
        )

    normalized: dict[str, dict[str, Any]] = {}
    for filename, metadata in documents.items():
        if not isinstance(filename, str) or not isinstance(metadata, dict):
            raise IngestionError(
                f"Manifesto inválido em '{manifest_path}': "
                "cada documento deve mapear para metadados"
            )
        normalized[filename] = dict(metadata)
    return normalized


def ingest_file(
    path: str | Path,
    *,
    metadata: Mapping[str, Any] | None = None,
) -> list[DocumentFragment]:
    """Lê um documento suportado e devolve seus fragmentos estruturais."""

    source_path = Path(path)
    if not source_path.exists():
        raise IngestionError(f"Documento não encontrado: {source_path}")
    if not source_path.is_file():
        raise IngestionError(f"O caminho não é um arquivo: {source_path}")

    extension = source_path.suffix.casefold()
    reader = _READERS.get(extension)
    if reader is None:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise UnsupportedDocumentError(
            f"Formato '{extension or '(sem extensão)'}' não suportado para '{source_path.name}'. "
            f"Formatos aceitos: {supported}"
        )

    document_metadata = dict(metadata or {})
    document_metadata.setdefault("format", extension.removeprefix("."))
    try:
        fragments = reader(source_path, document_metadata)
    except IngestionError:
        raise
    except Exception as exc:  # leitores de terceiros têm exceções específicas variadas
        raise IngestionError(f"Falha ao ler '{source_path.name}': {exc}") from exc

    fragments = [fragment for fragment in fragments if clean_text(fragment.text)]
    if not fragments:
        raise EmptyDocumentError(f"Nenhum conteúdo textual foi extraído de '{source_path.name}'")
    return fragments


def load_document(
    path: str | Path,
    *,
    metadata: Mapping[str, Any] | None = None,
) -> list[DocumentFragment]:
    """Alias legível para :func:`ingest_file`."""

    return ingest_file(path, metadata=metadata)


def load_documents(
    documents_dir: str | Path,
    manifest_path: str | Path | None = None,
) -> list[DocumentFragment]:
    """Carrega, em ordem estável, todos os documentos suportados de uma pasta.

    Arquivos de outros formatos são ignorados para permitir que a pasta contenha
    notas auxiliares. Subpastas não são percorridas implicitamente.
    """

    directory = Path(documents_dir)
    if not directory.exists():
        raise IngestionError(f"Pasta de documentos não encontrada: {directory}")
    if not directory.is_dir():
        raise IngestionError(f"O caminho não é uma pasta de documentos: {directory}")

    manifest = load_manifest(manifest_path)
    files = sorted(
        (
            item
            for item in directory.iterdir()
            if item.is_file() and item.suffix.casefold() in SUPPORTED_EXTENSIONS
        ),
        key=lambda item: item.name.casefold(),
    )

    fragments: list[DocumentFragment] = []
    for path in files:
        file_metadata = manifest.get(path.name, {})
        fragments.extend(ingest_file(path, metadata=file_metadata))
    return fragments


def ingest_directory(
    documents_dir: str | Path,
    *,
    manifest_path: str | Path | None = None,
) -> list[DocumentFragment]:
    """Alias compatível com chamadas que explicitam a operação de ingestão."""

    return load_documents(documents_dir, manifest_path)


def _base_fields(path: Path, metadata: Mapping[str, Any]) -> dict[str, Any]:
    copied_metadata = dict(metadata)
    return {
        "source_name": path.name,
        "source_path": path,
        "category": str(copied_metadata.get("category", "Geral")),
        "metadata": copied_metadata,
    }


def _fragment(
    path: Path,
    metadata: Mapping[str, Any],
    text: str,
    location: str,
    *,
    page: int | None = None,
    section: str | None = None,
    extra_metadata: Mapping[str, Any] | None = None,
) -> DocumentFragment | None:
    normalized = clean_text(text)
    if not normalized:
        return None
    fields = _base_fields(path, metadata)
    if extra_metadata:
        fields["metadata"].update(extra_metadata)
    return DocumentFragment(
        text=normalized,
        location=location,
        page=page,
        section=section,
        **fields,
    )


def _read_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            pass
    # Documentos legados em português frequentemente usam Windows-1252.
    return raw.decode("cp1252")


def _read_plain_text(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    fragment = _fragment(path, metadata, _read_text(path), "Documento")
    return [fragment] if fragment else []


_MARKDOWN_HEADING = re.compile(r"^(#{1,6})\s+(.+?)\s*#*\s*$")


def _split_markdown_sections(text: str) -> list[tuple[str | None, str]]:
    sections: list[tuple[str | None, str]] = []
    current_title: str | None = None
    current_lines: list[str] = []

    def flush() -> None:
        value = clean_text("\n".join(current_lines))
        if value:
            sections.append((current_title, value))

    for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        match = _MARKDOWN_HEADING.match(line.strip())
        if match:
            flush()
            current_title = clean_text(match.group(2))
            current_lines = [line]
        else:
            current_lines.append(line)
    flush()
    return sections


def _read_markdown(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    result: list[DocumentFragment] = []
    for section, text in _split_markdown_sections(_read_text(path)):
        location = f"Seção: {section}" if section else "Documento"
        fragment = _fragment(path, metadata, text, location, section=section)
        if fragment:
            result.append(fragment)
    return result


def _read_pdf(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - dependência declarada no projeto
        raise IngestionError("Leitura de PDF requer o pacote 'pypdf'") from exc

    reader = PdfReader(path)
    if reader.is_encrypted and reader.decrypt("") == 0:
        raise IngestionError(f"O PDF '{path.name}' está protegido por senha")

    result: list[DocumentFragment] = []
    for number, page in enumerate(reader.pages, start=1):
        fragment = _fragment(
            path,
            metadata,
            page.extract_text() or "",
            f"Página {number}",
            page=number,
            extra_metadata={"page": number},
        )
        if fragment:
            result.append(fragment)
    return result


def _read_docx(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    try:
        from docx import Document
        from docx.document import Document as DocumentType
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:  # pragma: no cover - dependência declarada no projeto
        raise IngestionError("Leitura de DOCX requer o pacote 'python-docx'") from exc

    document = Document(path)

    def blocks(parent: DocumentType) -> Iterable[Paragraph | Table]:
        for child in parent.element.body.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    result: list[DocumentFragment] = []
    section: str | None = None
    content: list[str] = []

    def flush() -> None:
        nonlocal content
        location = f"Seção: {section}" if section else "Documento"
        fragment = _fragment(path, metadata, "\n\n".join(content), location, section=section)
        if fragment:
            result.append(fragment)
        content = []

    for block in blocks(document):
        if isinstance(block, Paragraph):
            value = clean_text(block.text)
            if not value:
                continue
            style_name = (block.style.name if block.style else "").casefold()
            is_heading = any(
                marker in style_name for marker in ("heading", "título", "titulo", "title")
            )
            if is_heading:
                flush()
                section = value
            content.append(value)
        else:
            rows = [" | ".join(clean_text(cell.text) for cell in row.cells) for row in block.rows]
            table_text = clean_text("\n".join(rows))
            if table_text:
                content.append(table_text)
    flush()
    return result


def _cell_text(value: Any) -> str:
    if value is None:
        return ""
    if hasattr(value, "isoformat"):
        try:
            return str(value.isoformat())
        except TypeError:
            pass
    return str(value)


def _read_xlsx(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover - dependência declarada no projeto
        raise IngestionError("Leitura de XLSX requer o pacote 'openpyxl'") from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    result: list[DocumentFragment] = []
    try:
        for worksheet in workbook.worksheets:
            lines: list[str] = []
            for row in worksheet.iter_rows(values_only=True):
                values = [_cell_text(value) for value in row]
                while values and not values[-1]:
                    values.pop()
                if any(value.strip() for value in values):
                    lines.append(" | ".join(values))
            fragment = _fragment(
                path,
                metadata,
                "\n".join(lines),
                f"Planilha: {worksheet.title}",
                section=worksheet.title,
                extra_metadata={"sheet": worksheet.title},
            )
            if fragment:
                result.append(fragment)
    finally:
        workbook.close()
    return result


def _read_pptx(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - dependência declarada no projeto
        raise IngestionError("Leitura de PPTX requer o pacote 'python-pptx'") from exc

    presentation = Presentation(path)
    result: list[DocumentFragment] = []
    for number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        title: str | None = None
        if slide.shapes.title is not None:
            title = clean_text(slide.shapes.title.text)

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    value = clean_text(paragraph.text)
                    if value:
                        lines.append(value)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    value = " | ".join(clean_text(cell.text) for cell in row.cells)
                    if clean_text(value):
                        lines.append(value)

        fragment = _fragment(
            path,
            metadata,
            "\n".join(lines),
            f"Slide {number}",
            section=title,
            extra_metadata={"slide": number},
        )
        if fragment:
            result.append(fragment)
    return result


def _read_csv(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    raw = _read_text(path)
    sample = raw[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel

    rows = list(csv.reader(StringIO(raw), dialect=dialect))
    text = "\n".join(" | ".join(clean_text(cell) for cell in row) for row in rows)
    fragment = _fragment(path, metadata, text, "Tabela")
    return [fragment] if fragment else []


def _read_json(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    try:
        payload = json.loads(_read_text(path))
    except json.JSONDecodeError as exc:
        raise IngestionError(
            f"JSON inválido em '{path.name}', linha {exc.lineno}, coluna {exc.colno}: {exc.msg}"
        ) from exc

    result: list[DocumentFragment] = []
    if isinstance(payload, dict) and payload:
        for key, value in payload.items():
            section = str(key)
            rendered = json.dumps(value, ensure_ascii=False, indent=2, sort_keys=True)
            fragment = _fragment(
                path,
                metadata,
                f"{section}:\n{rendered}",
                f"Chave: {section}",
                section=section,
                extra_metadata={"json_key": section},
            )
            if fragment:
                result.append(fragment)
    else:
        rendered = json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True)
        fragment = _fragment(path, metadata, rendered, "Documento")
        if fragment:
            result.append(fragment)
    return result


def _read_html(path: Path, metadata: Mapping[str, Any]) -> list[DocumentFragment]:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:  # pragma: no cover - dependência declarada no projeto
        raise IngestionError("Leitura de HTML requer o pacote 'beautifulsoup4'") from exc

    soup = BeautifulSoup(_read_text(path), "html.parser")
    for unwanted in soup(["script", "style", "noscript", "template"]):
        unwanted.decompose()

    root = soup.body or soup
    selected_tags = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "blockquote", "tr"}
    result: list[DocumentFragment] = []
    section: str | None = None
    content: list[str] = []

    def flush() -> None:
        nonlocal content
        location = f"Seção: {section}" if section else "Documento"
        fragment = _fragment(path, metadata, "\n\n".join(content), location, section=section)
        if fragment:
            result.append(fragment)
        content = []

    for element in root.find_all(selected_tags):
        # Evita repetir texto quando um bloco selecionado está aninhado em outro.
        if element.find_parent(selected_tags):
            continue
        if element.name == "tr":
            value = " | ".join(
                cell.get_text(" ", strip=True) for cell in element.find_all(["th", "td"])
            )
        else:
            value = element.get_text(" ", strip=True)
        value = clean_text(value)
        if not value:
            continue
        if element.name in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            flush()
            section = value
        content.append(value)
    flush()

    if not result:
        fragment = _fragment(path, metadata, root.get_text("\n", strip=True), "Documento")
        if fragment:
            result.append(fragment)
    return result


_READERS: dict[str, Reader] = {
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
    ".md": _read_markdown,
    ".markdown": _read_markdown,
    ".txt": _read_plain_text,
    ".csv": _read_csv,
    ".json": _read_json,
    ".html": _read_html,
    ".htm": _read_html,
}


__all__ = [
    "EmptyDocumentError",
    "IngestionError",
    "SUPPORTED_EXTENSIONS",
    "UnsupportedDocumentError",
    "clean_text",
    "ingest_directory",
    "ingest_file",
    "load_document",
    "load_documents",
    "load_manifest",
]
