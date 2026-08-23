from __future__ import annotations

import json
from pathlib import Path

import pytest

from bimbam_agent.ingestion import (
    EmptyDocumentError,
    IngestionError,
    UnsupportedDocumentError,
    clean_text,
    ingest_file,
    load_documents,
    load_manifest,
)


def _write_pdf(path: Path, page_texts: list[str]) -> None:
    """Cria um PDF textual mínimo sem depender de serviços externos."""

    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_reference = writer._add_object(font)

    for text in page_texts:
        page = writer.add_blank_page(width=612, height=792)
        page[NameObject("/Resources")] = DictionaryObject(
            {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_reference})}
        )
        escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream = DecodedStreamObject()
        stream.set_data(f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1"))
        page[NameObject("/Contents")] = writer._add_object(stream)

    with path.open("wb") as output:
        writer.write(output)


def test_clean_text_is_conservative() -> None:
    raw = "\ufeff Título\r\n\r\n\r\n-  item\tcom   espaços\u00a0\r\n"
    assert clean_text(raw) == "Título\n\n- item com espaços"


def test_pdf_reader_preserves_page_and_manifest_metadata(tmp_path: Path) -> None:
    path = tmp_path / "manual.pdf"
    _write_pdf(path, ["Primeira pagina", "Segunda pagina"])

    fragments = ingest_file(
        path,
        metadata={"category": "Garantia", "language": "pt-BR"},
    )

    assert [fragment.page for fragment in fragments] == [1, 2]
    assert [fragment.location for fragment in fragments] == ["Página 1", "Página 2"]
    assert "Primeira pagina" in fragments[0].text
    assert fragments[0].category == "Garantia"
    assert fragments[0].metadata["language"] == "pt-BR"
    assert fragments[1].metadata["page"] == 2


def test_docx_reader_preserves_sections_and_table(tmp_path: Path) -> None:
    from docx import Document

    path = tmp_path / "manual.docx"
    document = Document()
    document.add_heading("Garantia", level=1)
    document.add_paragraph("Cobertura de doze meses.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Produto"
    table.cell(0, 1).text = "Prazo"
    table.cell(1, 0).text = "Notebook"
    table.cell(1, 1).text = "12 meses"
    document.save(path)

    fragments = ingest_file(path)

    assert len(fragments) == 1
    assert fragments[0].section == "Garantia"
    assert "Cobertura de doze meses." in fragments[0].text
    assert "Notebook | 12 meses" in fragments[0].text


def test_xlsx_reader_creates_one_fragment_per_nonempty_sheet(tmp_path: Path) -> None:
    from openpyxl import Workbook

    path = tmp_path / "prazos.xlsx"
    workbook = Workbook()
    active = workbook.active
    active.title = "Envios"
    active.append(["Região", "Dias"])
    active.append(["Sudeste", 3])
    empty = workbook.create_sheet("Vazia")
    empty["A1"] = None
    workbook.save(path)
    workbook.close()

    fragments = ingest_file(path)

    assert len(fragments) == 1
    assert fragments[0].section == "Envios"
    assert fragments[0].metadata["sheet"] == "Envios"
    assert "Sudeste | 3" in fragments[0].text


def test_pptx_reader_preserves_slide_number_and_title(tmp_path: Path) -> None:
    from pptx import Presentation

    path = tmp_path / "treinamento.pptx"
    presentation = Presentation()
    for title, body in (("Pagamentos", "Cartão e Pix"), ("Envios", "Rastreamento")):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    presentation.save(path)

    fragments = ingest_file(path)

    assert [fragment.location for fragment in fragments] == ["Slide 1", "Slide 2"]
    assert [fragment.section for fragment in fragments] == ["Pagamentos", "Envios"]
    assert fragments[0].metadata["slide"] == 1
    assert "Cartão e Pix" in fragments[0].text


def test_text_markdown_csv_json_and_html_readers(tmp_path: Path) -> None:
    files = {
        "nota.txt": "Linha simples",
        "guia.md": "# Pagamentos\nAceitamos Pix.\n\n## Cartões\nAceitamos Visa.",
        "tabela.csv": "método;prazo\nPix;imediato\n",
        "dados.json": json.dumps({"garantia": {"meses": 12}}, ensure_ascii=False),
        "pagina.html": (
            "<html><head><script>ignorar()</script></head><body>"
            "<h1>Devoluções</h1><p>Prazo de sete dias.</p></body></html>"
        ),
    }
    for name, content in files.items():
        (tmp_path / name).write_text(content, encoding="utf-8")

    text = ingest_file(tmp_path / "nota.txt")
    markdown = ingest_file(tmp_path / "guia.md")
    csv_fragments = ingest_file(tmp_path / "tabela.csv")
    json_fragments = ingest_file(tmp_path / "dados.json")
    html = ingest_file(tmp_path / "pagina.html")

    assert text[0].text == "Linha simples"
    assert [fragment.section for fragment in markdown] == ["Pagamentos", "Cartões"]
    assert "Pix | imediato" in csv_fragments[0].text
    assert json_fragments[0].section == "garantia"
    assert '"meses": 12' in json_fragments[0].text
    assert html[0].section == "Devoluções"
    assert "Prazo de sete dias." in html[0].text
    assert "ignorar" not in html[0].text


def test_load_documents_applies_optional_manifest_and_ignores_other_files(
    tmp_path: Path,
) -> None:
    documents = tmp_path / "documents"
    documents.mkdir()
    (documents / "b.txt").write_text("Segundo", encoding="utf-8")
    (documents / "a.txt").write_text("Primeiro", encoding="utf-8")
    (documents / "rascunho.bin").write_bytes(b"ignore")

    manifest_path = tmp_path / "manifest.json"
    manifest_path.write_text(
        json.dumps(
            {"documents": {"a.txt": {"category": "FAQ", "audience": "Atendimento"}}},
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    fragments = load_documents(documents, manifest_path)

    assert [fragment.source_name for fragment in fragments] == ["a.txt", "b.txt"]
    assert fragments[0].category == "FAQ"
    assert fragments[0].metadata["audience"] == "Atendimento"
    assert fragments[1].category == "Geral"


def test_clear_errors_for_invalid_inputs(tmp_path: Path) -> None:
    unsupported = tmp_path / "arquivo.xyz"
    unsupported.write_text("conteúdo", encoding="utf-8")
    empty = tmp_path / "vazio.txt"
    empty.write_text(" \n\n", encoding="utf-8")
    invalid_json = tmp_path / "quebrado.json"
    invalid_json.write_text('{"aberto": true', encoding="utf-8")
    invalid_manifest = tmp_path / "manifest.json"
    invalid_manifest.write_text("[]", encoding="utf-8")

    with pytest.raises(UnsupportedDocumentError, match="não suportado"):
        ingest_file(unsupported)
    with pytest.raises(EmptyDocumentError, match="Nenhum conteúdo"):
        ingest_file(empty)
    with pytest.raises(IngestionError, match="JSON inválido"):
        ingest_file(invalid_json)
    with pytest.raises(IngestionError, match="raiz deve ser um objeto"):
        load_manifest(invalid_manifest)
    with pytest.raises(IngestionError, match="não encontrado"):
        ingest_file(tmp_path / "ausente.txt")
